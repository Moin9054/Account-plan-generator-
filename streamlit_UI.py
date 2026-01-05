import streamlit as st
import json
import re
from io import BytesIO
from datetime import date

st.set_page_config(page_title="Account Plan Generator", page_icon="🤖", layout="wide")
st.title("🤖 Account Plan Generator")

SECTIONS = [
    "Account Overview",
    "Last Year Assessment",
    "Strategic Position Diagnosis",
    "Account Intelligence",
    "Internal Changes",
    "Growth Strategy",
    "Risks & Concerns",
    "Action Plan"
]

try:
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from fpdf import FPDF
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptxPt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor as PptxRGB
    
    try:
        from Chat_pipeline import llm_chat, retrieve_context
    except Exception as e:
        st.error(f"Pipeline Error: {e}")
        st.stop()
        
except ImportError as e:
    st.error(f"Missing Requirement: {e}")
    st.stop()

def insert_horizontal_line(doc):
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    pBdr = OxmlElement('w:pBdr')
    pPr.insert_element_before(pBdr, 'w:shd', 'w:tabs', 'w:suppressLineNumbers', 'w:wordWrap', 'w:overflowPunct', 'w:topLinePunct', 'w:autoSpaceDE', 'w:autoSpaceDN', 'w:bidi', 'w:adjustRightInd', 'w:snapToGrid', 'w:spacing', 'w:ind', 'w:contextualSpacing', 'w:mirrorIndents', 'w:suppressOverlap', 'w:jc', 'w:textDirection', 'w:textAlignment', 'w:textboxTightWrap', 'w:outlineLvl', 'w:divId', 'w:cnfStyle', 'w:rPr', 'w:sectPr', 'w:pPrChange')
    bottom = OxmlElement('w:bottom')
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), '6')
    bottom.set(qn('w:space'), '1')
    bottom.set(qn('w:color'), 'auto')
    pBdr.append(bottom)

def clear_placeholders(slide):
    for shape in list(slide.placeholders):
        if shape.placeholder_format.idx == 1:
            sp = shape.element
            sp.getparent().remove(sp)

def draw_process_flow(slide, steps, title_text):
    title = slide.shapes.title
    title.text = title_text
    clear_placeholders(slide)
    
    if not steps: return

    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(1.8)
    height = Inches(0.8)
    gap = Inches(0.4)
    
    for i, step in enumerate(steps):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.text = step.strip()
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxRGB(68, 114, 196)
        shape.text_frame.paragraphs[0].font.size = PptxPt(11)
        shape.text_frame.paragraphs[0].font.color.rgb = PptxRGB(255, 255, 255)
        
        if i < len(steps) - 1:
            arrow_left = left + width
            arrow_top = top + (height / 2) - Inches(0.1)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, gap, Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = PptxRGB(165, 165, 165)
        
        left += width + gap
        if left + width > Inches(9.5):
            left = Inches(0.5)
            top += height + Inches(0.5)

def draw_table_from_text(slide, title_text, table_lines):
    title = slide.shapes.title
    title.text = title_text
    clear_placeholders(slide)
    
    rows = []
    for line in table_lines:
        if "|---" in line or "|:--" in line: continue
        cells = [c.strip() for c in line.split('|') if c.strip() != '']
        if cells: rows.append(cells)
    
    if not rows: return
    
    rows_count = len(rows)
    cols_count = len(rows[0])
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.5 * rows_count)
    
    shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
    table = shape.table
    
    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_data in enumerate(row_data):
            if c_idx < len(rows[0]):
                cell = table.cell(r_idx, c_idx)
                cell.text = cell_data
                cell.text_frame.paragraphs[0].font.size = PptxPt(11)
                if r_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PptxRGB(68, 114, 196)
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.color.rgb = PptxRGB(255, 255, 255)

def create_docx(text):
    doc = Document()
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = title.add_run('Account Plan')
    run.font.name = 'Arial'
    run.font.size = Pt(30)
    run.font.color.rgb = RGBColor(0, 0, 0)
    run.bold = True
    
    lines = text.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if not line or line.startswith("###") or "|" in line or "FLOW:" in line:
            i += 1
            continue

        if line.startswith('# '):
            is_visual = False
            for k in range(i + 1, min(i + 5, len(lines))):
                if "|" in lines[k] or "FLOW:" in lines[k]:
                    is_visual = True
                    break
                if lines[k].startswith('#'): break
            
            if is_visual:
                i += 1
                continue

            insert_horizontal_line(doc)
            h = doc.add_heading('', level=1)
            r = h.add_run(line.replace('# ', '').upper())
            r.font.name = 'Arial'; r.font.size = Pt(17); r.font.color.rgb = RGBColor(0, 0, 0); r.font.bold = True

        elif line.startswith('## '):
            h = doc.add_heading('', level=2)
            r = h.add_run(line.replace('## ', ''))
            r.font.name = 'Arial'; r.font.size = Pt(14); r.font.color.rgb = RGBColor(0, 0, 0); r.font.bold = True

        elif line.startswith('* ') or line.startswith('- ') or line.startswith('+ '):
            p = doc.add_paragraph(style='List Bullet')
            clean = line[2:]
            if "**" in clean:
                parts = clean.split("**")
                for j, part in enumerate(parts):
                    r = p.add_run(part)
                    r.font.name = 'Arial'; r.font.size = Pt(12); r.font.color.rgb = RGBColor(0, 0, 0)
                    if j % 2 == 1: r.bold = True
            else:
                r = p.add_run(clean)
                r.font.name = 'Arial'; r.font.size = Pt(12); r.font.color.rgb = RGBColor(0, 0, 0)
        else:
            p = doc.add_paragraph()
            if "**" in line:
                parts = line.split("**")
                for j, part in enumerate(parts):
                    r = p.add_run(part)
                    r.font.name = 'Arial'; r.font.size = Pt(12); r.font.color.rgb = RGBColor(0, 0, 0)
                    if j % 2 == 1: r.bold = True
            else:
                r = p.add_run(line)
                r.font.name = 'Arial'; r.font.size = Pt(12); r.font.color.rgb = RGBColor(0, 0, 0)
        i += 1
        
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def create_pdf(text):
    class PDFObj(FPDF):
        def header(self): pass
        def footer(self):
            self.set_y(-15); self.set_font('Arial', 'I', 8); self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')
            
    pdf = PDFObj()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font('Arial', 'B', 24); pdf.cell(0, 10, 'Account Plan', 0, 1, 'L'); pdf.ln(10)

    lines = text.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if not line or "|" in line or "FLOW:" in line:
            i += 1
            continue
            
        clean = line.replace('**', '')
        if line.startswith('# '):
            is_visual = False
            for k in range(i + 1, min(i + 5, len(lines))):
                if "|" in lines[k] or "FLOW:" in lines[k]:
                    is_visual = True
                    break
                if lines[k].startswith('#'): break
            if is_visual:
                i += 1
                continue

            pdf.ln(5); pdf.set_draw_color(0, 0, 0); pdf.line(10, pdf.get_y(), 200, pdf.get_y()); pdf.ln(2)
            pdf.set_font("Arial", 'B', 14); pdf.set_text_color(0, 0, 0)
            pdf.multi_cell(0, 8, clean.replace('# ', '').upper())
        elif line.startswith('## '):
            pdf.ln(2); pdf.set_font("Arial", 'B', 12)
            pdf.multi_cell(0, 6, clean.replace('## ', ''))
        else:
            pdf.set_font("Arial", '', 11)
            pdf.multi_cell(0, 6, clean.lstrip('* -+'))
        i += 1
        
    return pdf.output(dest='S').encode('latin-1')

def create_pptx(text):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Strategic Account Plan"
    title_slide.placeholders[1].text = f"{st.session_state.account_data.get('Account Name', 'Client')}\n{date.today().strftime('%B %d, %Y')}"

    content_layout = prs.slides.add_slide(prs.slide_layouts[1]) 
    
    LEFT, TOP, WIDTH, HEIGHT = Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.5)
    
    curr_slide, text_frame, textbox = None, None, None
    curr_header, is_visual_content = "", False
    table_buf, table_mode = [], False
    
    def start_new_content_slide(title, is_continuation=False):
        nonlocal curr_slide, text_frame, textbox
        
        slide = prs.slides.add_slide(prs.slide_layouts[1]) 
        slide.shapes.title.text = title + (" (Cont.)" if is_continuation else "")
        
        for shape in list(slide.placeholders):
            if shape.placeholder_format.idx == 1:
                sp = shape.element
                sp.getparent().remove(sp)
                
        textbox = slide.shapes.add_textbox(LEFT, TOP, WIDTH, HEIGHT)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        return slide, text_frame, textbox

    lines = text.split('\n')
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        if not line or line == "###": 
            i += 1; continue
            
        clean = line.replace('**', '').replace('* ', '').strip()

        if "FLOW:" in line:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            raw_flow = line.split("FLOW:")[1]
            if not raw_flow.strip() and i + 1 < len(lines):
                raw_flow = lines[i+1]
                i += 1
            steps = [s.strip().replace('*', '') for s in re.split(r'->|-->|=>|→', raw_flow) if s.strip()]
            draw_process_flow(slide, steps, "Strategic Process Flow")
            curr_slide, text_frame, textbox = None, None, None 
            i += 1; continue
            
        if "|" in line:
            if not table_mode: table_mode = True; table_buf = []
            table_buf.append(line)
            i += 1; continue
            
        if table_mode and "|" not in line:
            if table_buf:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                draw_table_from_text(slide, curr_header + " (Data)", table_buf)
            table_mode = False; table_buf = []
            curr_slide, text_frame, textbox = None, None, None 

        if line.startswith('# '):
            curr_header = clean.replace('# ', '').upper()
            
            is_visual_content = False
            for k in range(i + 1, min(i + 5, len(lines))):
                if "|" in lines[k] or "FLOW:" in lines[k]:
                    is_visual_content = True
                    break
                if lines[k].startswith('#'): break

            if not is_visual_content:
                curr_slide, text_frame, textbox = start_new_content_slide(curr_header, is_continuation=False)
            else:
                curr_slide, text_frame, textbox = None, None, None
            i += 1; continue

        if curr_slide and text_frame:
            
            if len(text_frame.paragraphs) > 10:
                
                text_to_carry_over = line 

                curr_slide, text_frame, textbox = start_new_content_slide(curr_header, is_continuation=True)
                
                line_to_add = text_to_carry_over
            else:
                line_to_add = line
            
            match = re.match(r'\*?\s*\*\*(.*?)\*\*\s*:\s*(.*)', line_to_add)
            if match:
                label, content = match.group(1), match.group(2)
                p = text_frame.add_paragraph()
                p.text = label
                p.font.bold = True
                p.font.size = PptxPt(15) 
                p.space_before = PptxPt(8)
                
                for sent in re.split(r'(?<=[.!?]) +', content):
                    if len(sent.strip()) > 5:
                        sub = text_frame.add_paragraph()
                        sub.text = sent.strip()
                        sub.level = 1
                        sub.font.size = PptxPt(11) 
                        sub.space_before = PptxPt(1)
                
            elif line_to_add.startswith('## '):
                p = text_frame.add_paragraph()
                p.text = line_to_add.replace('## ', '')
                p.font.bold = True
                p.font.size = PptxPt(17)
                p.font.color.rgb = PptxRGB(68, 114, 196)
                p.space_before = PptxPt(12)
                
            elif line_to_add.startswith('* ') or line_to_add.startswith('- ') or line_to_add.startswith('+ '):
                p = text_frame.add_paragraph()
                p.text = clean
                p.level = 0
                p.font.size = PptxPt(12)
                p.space_before = PptxPt(4)
                
            else:
                p = text_frame.add_paragraph()
                p.text = clean
                p.level = 0
                p.font.size = PptxPt(12)
                p.space_before = PptxPt(4)
        
        i += 1
        
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def extract_smart_data(user_input):
    prompt = f"""
    Analyze input: "{user_input}"
    Task: Extract ANY factual account data mentioned (e.g. Revenue, Stakeholders, Competitors, Pain Points, Goals, Tech Stack).
    Output JSON ONLY. Example: {{"Revenue": "$10M", "Competitor": "Acme"}}
    If no new info, return {{}}
    """
    try:
        res = llm_chat([{"role": "user", "content": prompt}], temperature=0.0, json_mode=True)
        data = json.loads(re.search(r'\{.*\}', res, re.DOTALL).group(0))
        for k, v in data.items():
            if v and v not in ["None", "N/A"]:
                st.session_state.account_data[k] = v
    except: pass

def prune_questions(queue, user_input):
    if not queue: return []
    queue_str = "\n".join([f"{i}. {q}" for i, q in enumerate(queue)])
    prompt = f"""
    User said: "{user_input}"
    Future Questions:
    {queue_str}
    
    Task: Return JSON list of indices for questions that are now ANSWERED.
    Example: [0, 2]
    """
    try:
        res = llm_chat([{"role": "user", "content": prompt}], temperature=0.0, json_mode=True)
        indices = json.loads(re.search(r'\[.*\]', res, re.DOTALL).group(0))
        return [q for i, q in enumerate(queue) if i not in indices]
    except: return queue

def get_questions_for_section(section, level):
    context = retrieve_context(f"PLANNING PROMPTS {section} {level} HELP")
    prompt = f"""
    Context Chunk:
    {context}
    
    User Settings: Section: {section}, Level: {level}
    
    Task:
    1. Find the question block specifically for "{level}".
    2. Extract all bulleted questions.
    3. Output JSON list of strings.
    """
    response = llm_chat([{"role": "user", "content": prompt}], temperature=0.0, json_mode=True)
    try:
        return json.loads(re.search(r'\[.*\]', response, re.DOTALL).group(0))
    except:
        return []

def get_system_prompt(account_data, transcript_data):
    today_str = date.today().strftime("%B %d, %Y")
    help_level = account_data.get('Help Level', 'Comprehensive')
    return f"""
    You are a Senior Account Strategist. Generate a **Final Account Plan**.
    
    ### EXTRACTED DATA:
    {json.dumps(account_data, indent=2)}
    
    ### TRANSCRIPT:
    {transcript_data}
    
    ### METADATA:
    * **Client**: {account_data.get('Account Name')}
    * **Tier**: {account_data.get('Tier')}
    * **Date**: {today_str}
    * **Help Level**: {help_level}
    
    ### STRICT OUTPUT FORMAT (Markdown):
    **Client Company** : [Insert Name]
    **Account Name** : {account_data.get('Account Name')}
    **Planning Depth** : {help_level}
    **Date** : {today_str}
    
    # 1. ACCOUNT OVERVIEW
    * **Key Characteristics** : (Write 150-200 words analyzing market position).
    * **Primary Contact** : {account_data.get('Contact', 'Not Provided')} (Write 150 words on influence).
    * **Financial History** : (Write 150 words on 3-year trend).
    * **Past Key Projects** : (Write 150 words).
    * **Coming Year Revenue Target** : {account_data.get('Revenue', 'Extract from transcript')} (Write 150 words justifying target).
    
    # 2. LAST YEAR ASSESSMENT
    * **Completed Initiatives** : (Write 150 words).
    * **Missed Opportunities** : (Write 150 words).
    * **Strategic Hits** : (Write 150 words).
    * **Strategic Misses** : (Write 150 words).
    
    # 3. STRATEGIC POSITION DIAGNOSIS
    * **Current Position** : (Write 200 words).
    * **Growth Constraint** : (Write 200 words).
    
    # 4. ACCOUNT INTELLIGENCE
    * **Strategic Direction** : (Write 200 words).
    * **Leadership & Org** : (Write 200 words).
    
    # 5. INTERNAL CHANGES
    * **Internal Changes** : (Write 200 words).
    * **Opportunities** : (Write 200 words).
    
    # 6. GROWTH STRATEGY
    * **Next Beachhead Opportunity** : (Write 200 words).
    * **Proof Points Needed** : (Write 200 words).
    
    # 4. STAKEHOLDER & RELATIONSHIP MAP
    | Name | Role | Influence | Strategy |
    
    # 7. RISKS & CONCERNS
    * **Relationship Risks** : (Write 200 words).
    * **Competitive Risks** : {account_data.get('Competitors', 'Extract from transcript')} (Write 200 words).
    
    # 8. ACTION PLAN
    * **Full Year Key Actions** : (Write 200 words).
    * **Q1 Key Actions** : (Write 200 words).
    
    FLOW: Analyze Requirements -> Develop Strategy -> Present Proposal -> Close Deal
    """

if "messages" not in st.session_state:
    st.session_state.messages = [{"role": "assistant", "content": "Let's help you make an Account Plan, please tell the Account name."}]
if "account_data" not in st.session_state:
    st.session_state.account_data = {"Account Name": None, "Tier": None, "Help Level": None}
if "chat_stage" not in st.session_state: st.session_state.chat_stage = "awaiting_name" 
if "current_section_idx" not in st.session_state: st.session_state.current_section_idx = 0
if "plan_generated" not in st.session_state: st.session_state.plan_generated = False
if "section_history" not in st.session_state: st.session_state.section_history = {s: [] for s in SECTIONS}
if "question_queue" not in st.session_state: st.session_state.question_queue = []

for msg in st.session_state.messages:
    with st.chat_message(msg["role"]): st.markdown(msg["content"].replace("$", "\$"))

if st.session_state.plan_generated:
    st.divider(); st.subheader("📥 Download Report")
    plan = st.session_state.messages[-1]["content"]
    
    account_name = st.session_state.account_data.get("Account Name", "Account_Plan")
    base_filename = "".join(c for c in account_name if c.isalnum() or c in (' ', '_')).rstrip()
    
    c1, c2, c3 = st.columns(3)
    c1.download_button(
        "Word", 
        create_docx(plan), 
        f"{base_filename}.docx", 
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
    c2.download_button(
        "PDF", 
        create_pdf(plan), 
        f"{base_filename}.pdf", 
        "application/pdf"
    )
    c3.download_button(
        "PPT", 
        create_pptx(plan), 
        f"{base_filename}.pptx", 
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    st.divider()

if prompt := st.chat_input("Type your answer here"):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"): st.markdown(prompt.replace("$", "\$"))
    
    with st.spinner("Thinking..."):
        ai_response = None
        
        if st.session_state.chat_stage == "awaiting_name":
            val_prompt = f"""
            Analyze USER INPUT: "{prompt}"
            Task: Classify into EXACTLY one category.
            1. "social_query": The input is a question about your well-being (e.g., "How are you?", "Hope you are doing well?").
            2. "simple_greeting": The input is a general conversational opener (e.g., "Hi", "Hello", "Good morning") that is NOT a question.
            3. "name": A specific company or account name (e.g. "Microsoft", "Acme Corp").
            4. "gibberish": Random characters, single meaningless words like "asdf".
            
            Output JSON: {{"category": "social_query/simple_greeting/name/gibberish"}}
            """
            val_res = llm_chat([{"role": "user", "content": val_prompt}], temperature=0.0, json_mode=True)
            category = "gibberish"
            try: category = json.loads(re.search(r'\{.*\}', val_res, re.DOTALL).group(0)).get("category", "gibberish")
            except: pass
            
            if category == "social_query":
                ai_response = "I am fantastic! Let's help you make an account plan. What is the Account Name?"
            elif category == "simple_greeting":
                ai_response = "Hello! Lets get started on your account plan. What is the Account Name?"
            elif category == "gibberish":
                 ai_response = "I didn't recognize that as a valid account name. Please enter the Client/Company Name."
            else:
                st.session_state.account_data["Account Name"] = prompt
                st.session_state.chat_stage = "awaiting_level"
                ai_response = f"""
                Great. Let's plan for **{prompt}**. Please select a help level:
                
                * **Streamlined**: Basic planning (2-3 questions).
                * **Guided**: Balanced planning (4-5 questions).
                * **Comprehensive**: In-depth planning (6-8 questions).
                """

        elif st.session_state.chat_stage == "awaiting_level":
            p_low = prompt.lower()
            if any(x in p_low for x in ["stream", "guided", "comp"]):
                lvl = "Streamlined" if "stream" in p_low else "Guided" if "guided" in p_low else "Comprehensive"
                st.session_state.account_data["Help Level"] = lvl
                st.session_state.chat_stage = "interview"
                
                st.session_state.current_section_idx = 0
                curr_sec = SECTIONS[0]
                q_list = get_questions_for_section(curr_sec, lvl)
                
                if not q_list:
                    q_list = [f"Please describe the **{curr_sec}** for this account."]
                
                st.session_state.question_queue = q_list
                q1 = st.session_state.question_queue.pop(0)
                ai_response = f"Great choice. Let's build a **{lvl}** plan.\n\n**Starting Phase: {curr_sec}**\n\n{q1}"
                
            else:
                ai_response = "Please select: Streamlined, Guided, or Comprehensive."

        elif st.session_state.chat_stage == "interview":
            
            val_prompt = f"""
            Analyze INPUT: "{prompt}"
            Classify: 
            1. "stop": User says "that's all", "done", "generate", "no more info".
            2. "greeting": Hi, Hello, how are you.
            3. "valid": Anything else.
            Output JSON: {{"status": "..."}}
            """
            val_res = llm_chat([{"role": "user", "content": val_prompt}], temperature=0.0, json_mode=True)
            status = "valid"
            try: status = json.loads(re.search(r'\{.*\}', val_res, re.DOTALL).group(0)).get("status", "valid")
            except: pass
            
            if status == "stop":
                st.session_state.chat_stage = "generating"
                
            elif status == "greeting":
                curr = SECTIONS[st.session_state.current_section_idx]
                ai_response = f"I am doing great! Let's get back to **{curr}**."
            
            else:
                extract_smart_data(prompt)
                
                curr_sec = SECTIONS[st.session_state.current_section_idx]
                st.session_state.section_history[curr_sec].append(f"Answer: {prompt}")
                
                if st.session_state.question_queue:
                    st.session_state.question_queue = prune_questions(st.session_state.question_queue, prompt)
                
                ack_prompt = f"User Answer: '{prompt}'. Write short 5-word acknowledgement. No 'That makes sense'."
                ack = llm_chat([{"role": "user", "content": ack_prompt}], temperature=0.7)
                
                if st.session_state.question_queue:
                    next_q = st.session_state.question_queue.pop(0)
                    ai_response = f"{ack}\n\n{next_q}"
                else:
                    st.session_state.current_section_idx += 1
                    found_next = False
                    
                    while st.session_state.current_section_idx < len(SECTIONS):
                        next_sec = SECTIONS[st.session_state.current_section_idx]
                        lvl = st.session_state.account_data.get('Help Level')

                        if "Strategic Position" in next_sec and lvl == "Streamlined":
                            st.session_state.current_section_idx += 1
                            continue
                        
                        new_qs = get_questions_for_section(next_sec, lvl)

                        if not new_qs:
                            new_qs = [f"Please provide details on **{next_sec}**."]
                        
                        st.session_state.question_queue = new_qs
                        q1 = st.session_state.question_queue.pop(0)
                        ai_response = f"{ack}\n\n**Moving on to: {next_sec}**\n\n{q1}"
                        found_next = True
                        break
                    
                    if not found_next:
                        st.session_state.chat_stage = "generating"

        if st.session_state.chat_stage == "generating" and not ai_response:
            transcript_text = ""
            for sec, items in st.session_state.section_history.items():
                transcript_text += f"SECTION: {sec}\n" + "\n".join(items) + "\n\n"
            
            final_prompt = get_system_prompt(st.session_state.account_data, transcript_text)
            
            with st.status("Generating Report"):
                ai_response = llm_chat([{"role": "user", "content": final_prompt}], temperature=0.3)
                if ai_response: st.session_state.plan_generated = True
                else: ai_response = "Error: Generation failed."

        if ai_response:
            st.session_state.messages.append({"role": "assistant", "content": ai_response})
        else:
            if st.session_state.chat_stage != "generating": st.error("Connection Error.")
            
        st.rerun()